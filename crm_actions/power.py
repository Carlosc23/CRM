from pptx import Presentation
from pptx.util import Inches

def generate_power():
    img_path = 'Graphic1.png'

    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Resumen de datos al día de hoy"
    subtitle.text = "Hospital"


    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)

    left = top = Inches(1)


    left = Inches(1)
    height = Inches(5.5)
    pic = slide.shapes.add_picture(img_path, left, top, height=height)
    slide = prs.slides.add_slide(blank_slide_layout)
    img_path = 'Graphic2.png'
    pic = slide.shapes.add_picture(img_path, left, top, height=height)
    slide = prs.slides.add_slide(blank_slide_layout)
    img_path = 'Graphic3.png'
    pic = slide.shapes.add_picture(img_path, left, top, height=height)
    slide = prs.slides.add_slide(blank_slide_layout)
    img_path = 'Graphic4.png'
    pic = slide.shapes.add_picture(img_path, left, top, height=height)
    slide = prs.slides.add_slide(blank_slide_layout)
    img_path = 'Graphic5.png'
    pic = slide.shapes.add_picture(img_path, left, top, height=height)
    slide = prs.slides.add_slide(blank_slide_layout)
    img_path = 'Graphic6.png'
    pic = slide.shapes.add_picture(img_path, left, top, height=height)
    prs.save('../static/test.pptx')